from __future__ import annotations

import argparse
import asyncio
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import wave
from pathlib import Path

import numpy as np
from pdf2image import convert_from_path
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SPEAKER_LINE_RE = re.compile(r"^\s*Speaker\s+([AB])\s*:\s*(.+?)\s*$", re.I)
PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

_TTS_RUNTIME: dict | None = None


def _run_cmd(cmd: list[str]) -> bytes:
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
    if result.returncode != 0:
        raise RuntimeError(f"{' '.join(cmd)}\n{result.stderr.decode(errors='ignore')}")
    return result.stdout


def _resolve_libreoffice_cmd() -> list[str]:
    env_override = os.getenv("LIBREOFFICE_CMD") or os.getenv("SOFFICE_CMD")
    if env_override:
        return [env_override]
    libreoffice_path = shutil.which("libreoffice")
    if libreoffice_path:
        return [libreoffice_path]
    soffice_path = shutil.which("soffice")
    if soffice_path:
        return [soffice_path]
    raise FileNotFoundError(
        "Neither 'libreoffice' nor 'soffice' was found in PATH. "
        "Set LIBREOFFICE_CMD to the full executable path if needed."
    )


def _get_tts_runtime() -> dict:
    global _TTS_RUNTIME
    if _TTS_RUNTIME is not None:
        return _TTS_RUNTIME

    megatts_root = Path(__file__).resolve().parents[1] / "presentagent" / "MegaTTS3"
    if str(megatts_root) not in sys.path:
        sys.path.append(str(megatts_root))

    from presentagent.MegaTTS3.tts.infer_cli import MegaTTS3DiTInfer

    ckpt_root = str(megatts_root / "checkpoints")
    infer = MegaTTS3DiTInfer(ckpt_root=ckpt_root)
    assets_dir = megatts_root / "assets"
    default_speaker = os.getenv("TTS_DEFAULT_SPEAKER", "B").strip().upper() or "B"
    if default_speaker not in {"A", "B"}:
        default_speaker = "B"
    speaker_a_prompt = Path(
        os.getenv(
            "TTS_SPEAKER_A_PROMPT",
            str(assets_dir / "00014.wav"),
        )
    )
    speaker_b_prompt = Path(
        os.getenv(
            "TTS_SPEAKER_B_PROMPT",
            str(assets_dir / "00009.wav"),
        )
    )

    def build_resource_context(prompt_audio_path: Path):
        audio_bytes = prompt_audio_path.read_bytes()
        latent_file = None
        potential_npy = prompt_audio_path.with_suffix(".npy")
        if potential_npy.is_file():
            latent_file = str(potential_npy)
        return infer.preprocess(audio_bytes, latent_file)

    speaker_contexts = {
        "A": build_resource_context(speaker_a_prompt),
        "B": build_resource_context(speaker_b_prompt),
    }
    _TTS_RUNTIME = {
        "infer": infer,
        "default_speaker": default_speaker,
        "speaker_contexts": speaker_contexts,
    }
    return _TTS_RUNTIME


def _generate_tts_audio(text: str, output_path: str) -> None:
    try:
        from presentagent.MegaTTS3.tts.utils.audio_utils.io import save_wav

        runtime = _get_tts_runtime()
        infer = runtime["infer"]
        default_speaker = runtime["default_speaker"]
        speaker_contexts = runtime["speaker_contexts"]

        def parse_speaker_segments(script: str) -> list[tuple[str, str]]:
            segments: list[tuple[str, str]] = []
            current_speaker = default_speaker
            buffer: list[str] = []
            for raw_line in script.splitlines():
                line = raw_line.strip()
                if not line:
                    continue
                match = SPEAKER_LINE_RE.match(line)
                if match:
                    if buffer:
                        segments.append((current_speaker, " ".join(buffer).strip()))
                        buffer = []
                    current_speaker = match.group(1).upper()
                    content = match.group(2).strip()
                    if content:
                        buffer.append(content)
                else:
                    buffer.append(line)
            if buffer:
                segments.append((current_speaker, " ".join(buffer).strip()))
            if not segments:
                segments.append((default_speaker, script.strip()))
            return [(speaker, content) for speaker, content in segments if content]

        segments = parse_speaker_segments(text)
        with tempfile.TemporaryDirectory() as temp_dir:
            segment_paths: list[str] = []
            for idx, (speaker, content) in enumerate(segments):
                wav_bytes = infer.forward(
                    speaker_contexts[speaker],
                    content,
                    time_step=32,
                    p_w=1.6,
                    t_w=2.5,
                )
                segment_path = os.path.join(temp_dir, f"speaker_{speaker}_{idx}.wav")
                save_wav(wav_bytes, segment_path)
                segment_paths.append(segment_path)
            if len(segment_paths) == 1:
                Path(output_path).write_bytes(Path(segment_paths[0]).read_bytes())
                return
            list_file_path = os.path.join(temp_dir, "segments.txt")
            with open(list_file_path, "w", encoding="utf-8") as f:
                for seg in segment_paths:
                    f.write(f"file '{seg}'\n")
            _run_cmd(
                [
                    "ffmpeg",
                    "-y",
                    "-f",
                    "concat",
                    "-safe",
                    "0",
                    "-i",
                    list_file_path,
                    "-c",
                    "copy",
                    output_path,
                ]
            )
    except Exception:
        print("TTS generation failed; falling back to 3-second silent audio.", file=sys.stderr)
        traceback.print_exc()
        if os.getenv("PPTAGENT_TTS_STRICT", "0") == "1":
            raise
        sample_rate = 22050
        duration = 3.0
        samples = np.zeros(int(sample_rate * duration), dtype=np.int16)
        with wave.open(output_path, "w") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(samples.tobytes())


def _create_video_segment(image_path: str, audio_path: str, temp_dir: str, index: int) -> str:
    output_path = os.path.join(temp_dir, f"segment_{index}.mp4")
    _run_cmd(
        [
            "ffmpeg",
            "-y",
            "-loop",
            "1",
            "-i",
            image_path,
            "-i",
            audio_path,
            "-vf",
            "scale=1920:1080",
            "-c:v",
            "libx264",
            "-tune",
            "stillimage",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            output_path,
        ]
    )
    return output_path


def _extract_slide_video_overlays(slide, temp_dir: str, slide_index: int) -> list[dict]:
    overlays: list[dict] = []
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.MEDIA:
            continue

        media_rid = None
        media_node = shape._element.find(".//p14:media", namespaces=PPTX_NS)
        if media_node is not None:
            media_rid = media_node.get(f"{{{PPTX_NS['r']}}}embed")
        if media_rid is None:
            video_node = shape._element.find(".//a:videoFile", namespaces=PPTX_NS)
            if video_node is not None:
                media_rid = video_node.get(f"{{{PPTX_NS['r']}}}link")
        if media_rid is None or media_rid not in slide.part.rels:
            continue

        rel = slide.part.rels[media_rid]
        target_part = getattr(rel, "target_part", None)
        blob = getattr(target_part, "blob", None)
        if blob is None:
            continue

        partname = os.path.basename(str(getattr(target_part, "partname", "media.bin")))
        overlay_path = os.path.join(
            temp_dir, f"slide_{slide_index}_shape_{shape_index}_{partname}"
        )
        with open(overlay_path, "wb") as f:
            f.write(blob)

        overlays.append(
            {
                "path": overlay_path,
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            }
        )
    return overlays


def _create_video_segment_with_overlays(
    image_path: str,
    audio_path: str,
    temp_dir: str,
    index: int,
    slide_width: int,
    slide_height: int,
    video_overlays: list[dict],
) -> str:
    output_path = os.path.join(temp_dir, f"segment_{index}.mp4")
    command = ["ffmpeg", "-y", "-loop", "1", "-i", image_path]
    for overlay in video_overlays:
        command.extend(["-stream_loop", "-1", "-i", overlay["path"]])
    command.extend(["-i", audio_path])

    overlay_speed = float(os.getenv("PPTAGENT_OVERLAY_SPEED", "1.0") or "1.0")
    if overlay_speed <= 0:
        overlay_speed = 1.0
    overlay_pts_factor = 1.0 / overlay_speed

    filter_parts = ["[0:v]scale=1920:1080[bg0]"]
    current_label = "bg0"
    for idx, overlay in enumerate(video_overlays, start=1):
        scaled_width = max(2, round(overlay["width"] * 1920 / slide_width))
        scaled_height = max(2, round(overlay["height"] * 1080 / slide_height))
        scaled_width += scaled_width % 2
        scaled_height += scaled_height % 2
        x = max(0, round(overlay["left"] * 1920 / slide_width))
        y = max(0, round(overlay["top"] * 1080 / slide_height))
        overlay_label = f"ov{idx}"
        next_label = f"bg{idx}"
        filter_parts.append(
            f"[{idx}:v]setpts={overlay_pts_factor:.6f}*PTS,scale={scaled_width}:{scaled_height}[{overlay_label}]"
        )
        filter_parts.append(
            f"[{current_label}][{overlay_label}]overlay={x}:{y}:shortest=1[{next_label}]"
        )
        current_label = next_label

    audio_input_index = len(video_overlays) + 1
    command.extend(
        [
            "-filter_complex",
            ";".join(filter_parts),
            "-map",
            f"[{current_label}]",
            "-map",
            f"{audio_input_index}:a",
            "-c:v",
            "libx264",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            output_path,
        ]
    )
    _run_cmd(command)
    return output_path


def _merge_video_segments(video_segments: list[str], output_path: str) -> None:
    list_file_path = output_path.replace(".mp4", "_list.txt")
    with open(list_file_path, "w", encoding="utf-8") as f:
        for seg in video_segments:
            f.write(f"file '{seg}'\n")
    try:
        _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-f",
                "concat",
                "-safe",
                "0",
                "-i",
                list_file_path,
                "-c",
                "copy",
                output_path,
            ]
        )
    finally:
        if os.path.exists(list_file_path):
            os.remove(list_file_path)


async def _build_video(ppt_path: Path, output_dir: Path) -> None:
    pdf_path = output_dir / f"{ppt_path.stem}.pdf"
    libreoffice_cmd = _resolve_libreoffice_cmd()
    _run_cmd(
        libreoffice_cmd + [
            "--headless",
            "--convert-to",
            "pdf",
            str(ppt_path),
            "--outdir",
            str(output_dir),
        ]
    )

    images_from_path = await asyncio.to_thread(convert_from_path, str(pdf_path))
    prs = await asyncio.to_thread(PptxPresentation, str(ppt_path))
    if len(images_from_path) != len(prs.slides):
        raise RuntimeError("Slide count does not match rendered image count.")

    with tempfile.TemporaryDirectory() as temp_dir:
        notes_dir = output_dir / "notes_assets"
        notes_dir.mkdir(parents=True, exist_ok=True)
        video_segments: list[str] = []

        notes_dump = []
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            if not notes.strip():
                notes = f"This is slide {i + 1}"
            notes_dump.append({"slide_idx": i + 1, "notes": notes})

            image_path = notes_dir / f"frame_{i}.jpg"
            audio_path = notes_dir / f"frame_{i}.wav"
            image.save(image_path)
            await asyncio.to_thread(_generate_tts_audio, notes, str(audio_path))
            video_overlays = _extract_slide_video_overlays(slide, str(notes_dir), i)
            if video_overlays:
                segment_path = await asyncio.to_thread(
                    _create_video_segment_with_overlays,
                    str(image_path),
                    str(audio_path),
                    str(notes_dir),
                    i,
                    prs.slide_width,
                    prs.slide_height,
                    video_overlays,
                )
            else:
                segment_path = await asyncio.to_thread(
                    _create_video_segment,
                    str(image_path),
                    str(audio_path),
                    str(notes_dir),
                    i,
                )
            video_segments.append(segment_path)

        (output_dir / "slide_notes.json").write_text(
            json.dumps(notes_dump, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        final_video_path = output_dir / "output.mp4"
        await asyncio.to_thread(_merge_video_segments, video_segments, str(final_video_path))

    print("success: True")
    print(f"pptx: {ppt_path}")
    print(f"pdf: {pdf_path}")
    print(f"video: {output_dir / 'output.mp4'}")
    print(f"notes: {output_dir / 'slide_notes.json'}")
    print(f"assets_dir: {output_dir / 'notes_assets'}")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Render a PPT with notes into a narrated MP4 and save intermediate assets."
    )
    parser.add_argument("--pptx", required=True, help="Path to final.pptx")
    parser.add_argument("--output-dir", required=True, help="Directory to save video outputs")
    args = parser.parse_args()

    ppt_path = Path(args.pptx).resolve()
    if not ppt_path.exists():
        raise SystemExit(f"PPTX not found: {ppt_path}")
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    asyncio.run(_build_video(ppt_path, output_dir))


if __name__ == "__main__":
    main()
