from __future__ import annotations

import base64
import io
import re
import zipfile
from pathlib import Path
import xml.etree.ElementTree as ET

from PIL import Image
from pptx import Presentation


PNG_DATA_RE = re.compile(r"data:image/png;base64,([A-Za-z0-9+/=\s]+)")


def extract_png_from_svg(svg_path: Path, out_path: Path) -> Path:
    text = svg_path.read_text(encoding="utf-8", errors="ignore")
    match = PNG_DATA_RE.search(text)
    if not match:
        raise RuntimeError(f"No embedded PNG found in {svg_path}")
    png_bytes = base64.b64decode(re.sub(r"\s+", "", match.group(1)))
    out_path.write_bytes(png_bytes)
    return out_path


def replace_slide_text(slide, title_text: str, body_lines: list[str]) -> None:
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        title_shape.text = title_text

    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape and shape.text_frame is not None:
            if getattr(shape, "is_placeholder", False) or shape.text.strip():
                body_shape = shape
                break
    if body_shape is None:
        raise RuntimeError("Could not find body text shape to replace.")

    tf = body_shape.text_frame
    tf.clear()
    for idx, line in enumerate(body_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line


def replace_primary_picture(slide, image_path: Path) -> None:
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    if not picture_shapes:
        raise RuntimeError("No picture placeholder found on slide.")
    target = picture_shapes[0]
    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    sample_root = repo_root / "robot_vla_worldmodel_batch15" / "EgoScale"
    ppt_path = sample_root / "document_to_ppt" / "single_presentation" / "final_single_presentation.pptx"
    assets_dir = sample_root / "url_to_source" / "assets"
    generated_dir = sample_root / "document_to_ppt" / "single_presentation" / "manual_fix_assets"
    generated_dir.mkdir(parents=True, exist_ok=True)

    pipeline_png = extract_png_from_svg(
        assets_dir / "4b7b2cb4_Pipeline.svg",
        generated_dir / "pipeline_embedded.png",
    )
    model_png = extract_png_from_svg(
        assets_dir / "3eb10577_Model_arch.svg",
        generated_dir / "model_arch_embedded.png",
    )

    prs = Presentation(str(ppt_path))

    slide3 = prs.slides[2]
    replace_primary_picture(slide3, pipeline_png)
    replace_slide_text(
        slide3,
        "EgoScale Pipeline",
        [
            "1. Pretrain a flow-based VLA on 20,854 hours of egocentric human videos.",
            "2. Mid-train with aligned human-robot play to bridge sensing and control.",
            "3. Post-train on downstream robot tasks for dexterous adaptation and transfer.",
        ],
    )

    slide4 = prs.slides[3]
    replace_primary_picture(slide4, model_png)
    replace_slide_text(
        slide4,
        "Flow-Based VLA Policy",
        [
            "1. A VLM backbone encodes visual observations and language instructions.",
            "2. A DiT action expert predicts wrist-level actions in a unified representation.",
            "3. The shared action space lets human demonstrations transfer to robot control.",
        ],
    )

    backup_path = ppt_path.with_name("final_single_presentation.before_egoscale_fix.pptx")
    if not backup_path.exists():
        backup_path.write_bytes(ppt_path.read_bytes())
    prs.save(str(ppt_path))

    with zipfile.ZipFile(ppt_path, "r") as zf:
        media_names = [name for name in zf.namelist() if name.startswith("ppt/media/")]
    manifest = {
        "ppt_path": str(ppt_path),
        "backup_path": str(backup_path),
        "generated_images": [str(pipeline_png), str(model_png)],
        "ppt_media": media_names,
    }
    (generated_dir / "fix_manifest.json").write_text(
        __import__("json").dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(ppt_path)


if __name__ == "__main__":
    main()
